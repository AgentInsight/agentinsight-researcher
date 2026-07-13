"""API 路由: OpenAI 兼容端点 + 文件上传.

OpenAI 兼容端点与 API 测试硬约束:
- 统一调用 OpenAI 兼容端点 POST /v1/chat/completions, 请求体带 stream: true
- 测试页面只能走对外 OpenAI 兼容接口, 禁止调用后端私有端点
- API 测试必须覆盖流式 SSE + 非流式 + 错误码
- 包含携带 Bearer JWT Token 与不携带两种场景
"""

from __future__ import annotations

import asyncio
import logging
import time
import uuid
from pathlib import Path
from typing import Any

import orjson
from fastapi import (
    APIRouter,
    File,
    Header,
    HTTPException,
    Request,
    UploadFile,
)
from fastapi.responses import JSONResponse, Response, StreamingResponse
from langchain_core.messages import AIMessage, BaseMessage, HumanMessage
from pydantic import BaseModel, Field, StrictBool

from src.api.middleware import (
    get_request_agent_id,
    get_request_session_id,
    get_request_user_id,
)
from src.config.settings import get_settings
from src.observability.tracing import trace_agent
from src.skills.researcher.chitchat_responder import get_chitchat_responder
from src.skills.researcher.query_classifier import (
    QueryIntent,
    get_query_intent_classifier,
)

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/v1", tags=["openai-compatible"])

# 后台任务引用保留 (防止 GC 静默取消 asyncio.create_task)
_background_tasks: set[asyncio.Task[Any]] = set()


def _create_background_task(coro: Any) -> asyncio.Task[Any]:
    """创建后台任务并保留引用 (防止 GC 静默取消).

    标准模式: set + done_callback(discard), 任务完成后自动从集合移除.
    """
    task = asyncio.create_task(coro)
    _background_tasks.add(task)
    task.add_done_callback(_background_tasks.discard)
    return task


# ========== LangGraph 图单例 (延迟构建, 复用) ==========
_compiled_graph: Any | None = None
_multi_agent_graph: Any | None = None
_chat_graph: Any | None = None  # 对话追问图单例


async def _get_graph(multi_agent: bool = False) -> Any:
    """获取/构建已编译的 LangGraph 单例.

    生产 StateGraph 必须挂 PostgresSaver.
    首次调用时构建, 后续复用 (单例).
    multi_agent=True 时构建多 Agent Supervisor 图.
    """
    global _compiled_graph, _multi_agent_graph
    if multi_agent:
        if _multi_agent_graph is None:
            from src.graph.multi_agent_builder import build_multi_agent_graph

            _multi_agent_graph = await build_multi_agent_graph(use_checkpointer=True)
        return _multi_agent_graph
    if _compiled_graph is None:
        from src.graph.builder import build_researcher_graph

        _compiled_graph = await build_researcher_graph(use_checkpointer=True)
    return _compiled_graph


async def _get_chat_graph() -> Any:
    """获取/构建对话追问图单例.

    生产 StateGraph 必须挂 PostgresSaver.
    复用同一 thread_id 隔离, 支持多会话并发.
    单节点 chat 图, 依赖 checkpointer 自动加载会话历史 (report_md / messages).
    """
    global _chat_graph
    if _chat_graph is None:
        from src.graph.chat_builder import build_chat_graph

        _chat_graph = await build_chat_graph(use_checkpointer=True)
    return _chat_graph


async def _has_report(session_id: str) -> bool:
    """检查会话是否已有报告.

    优先查 research_reports 表 (有 session_id 索引), 避免每次
    aget_state 加载全量 State; 查询失败时降级回 aget_state.

    thread_id 做会话隔离, checkpointer 自动持久化.

    Args:
        session_id: 会话 ID (thread_id)

    Returns:
        True 表示会话已有报告 (用于意图分类 has_report 参数).
    """
    # 优先走 report_store (按 session_id 索引查询, 比 aget_state 快几个数量级)
    try:
        from src.memory.report_store import get_report_store

        report_store = get_report_store()
        # list_reports 已按 session_id 过滤, 取 1 条即可判断
        reports = await report_store.list_reports(
            user_id=None,  # _has_report 仅按 session_id 判断, 不区分 user_id
            session_id=session_id,
            limit=1,
        )
        if reports:
            return True
    except Exception as e:  # noqa: BLE001
        logger.warning("report_store 查询失败, 降级 aget_state (session=%s): %s", session_id, e)

    # report_store 查询失败时回退到 aget_state
    try:
        graph = await _get_graph()
        config: dict[str, Any] = {"configurable": {"thread_id": session_id}}
        state_snapshot = await graph.aget_state(config)
        if state_snapshot and state_snapshot.values:
            report_md = state_snapshot.values.get("report_md", "")
            return bool(report_md)
    except Exception as e:  # noqa: BLE001
        logger.warning("检查会话报告失败 (session=%s): %s", session_id, e)
    return False


async def _load_chitchat_history(session_id: str) -> list[dict[str, str]]:
    """从 chat graph checkpointer 加载对话历史 (会话持久化).

    会话级数据通过 Checkpointer 隔离, thread_id 从请求上下文注入.
    ChitchatResponder 绕过 LangGraph 图, 需手动从 checkpointer 读取历史消息,
    否则跨请求上下文丢失.

    Args:
        session_id: 会话 ID (thread_id)

    Returns:
        历史消息列表 [{"role": "user"/"assistant", "content": "..."}]
    """
    try:
        graph = await _get_chat_graph()
        config: dict[str, Any] = {"configurable": {"thread_id": session_id}}
        state_snapshot = await graph.aget_state(config)
        if not state_snapshot or not state_snapshot.values:
            return []
        messages: list[BaseMessage] = state_snapshot.values.get("messages", []) or []
        history: list[dict[str, str]] = []
        for msg in messages:
            if isinstance(msg, HumanMessage):
                content = msg.content if isinstance(msg.content, str) else str(msg.content)
                history.append({"role": "user", "content": content})
            elif isinstance(msg, AIMessage):
                content = msg.content if isinstance(msg.content, str) else str(msg.content)
                history.append({"role": "assistant", "content": content})
        return history
    except Exception as e:  # noqa: BLE001
        logger.warning("加载闲聊历史失败 (session=%s): %s", session_id, e)
        return []


async def _save_chitchat_response(
    session_id: str,
    query: str,
    response: str,
    *,
    user_id: str | None = None,
    agent_id: str | None = None,
) -> None:
    """将 ChitchatResponder 的响应保存回 chat graph checkpointer (会话持久化).

    会话持久化到 Postgres Checkpointer.
    ChitchatResponder 绕过 LangGraph 图, 需手动写入 checkpointer,
    否则下一次请求无法读到本次对话上下文.

    使用 graph.aupdate_state 增量写入 messages (add_messages reducer 自动合并),
    不触发节点执行 (as_node=None).

    Args:
        session_id: 会话 ID (thread_id)
        query: 用户查询 (保存为 HumanMessage)
        response: ChitchatResponder 响应 (保存为 AIMessage)
        user_id: 用户 ID (可选, 写入 state 元数据)
        agent_id: Agent ID (可选, 写入 state 元数据)
    """
    try:
        graph = await _get_chat_graph()
        config: dict[str, Any] = {"configurable": {"thread_id": session_id}}
        await graph.aupdate_state(
            config,
            {
                "messages": [HumanMessage(content=query), AIMessage(content=response)],
                "session_id": session_id,
                "user_id": user_id,
                "agent_id": agent_id,
            },
        )
    except Exception as e:  # noqa: BLE001
        logger.warning("保存闲聊响应失败 (session=%s): %s", session_id, e)


# ========== 请求/响应模型 (OpenAI 兼容) ==========


class ChatMessage(BaseModel):
    """OpenAI 兼容消息格式."""

    role: str = "user"
    content: str = ""


class ChatCompletionRequest(BaseModel):
    """OpenAI 兼容 /v1/chat/completions 请求."""

    model: str = "agentinsight-researcher"
    messages: list[ChatMessage]
    stream: StrictBool = False
    temperature: float | None = None
    max_tokens: int | None = None
    # 扩展字段 (非 OpenAI 标准, 用于研究配置)
    report_type: str | None = Field(
        None, description="报告类型: basic_report | detailed_report | deep_research"
    )
    report_format: str | None = Field(None, description="输出格式: markdown | html | pdf")
    tone: str | None = Field(
        None, description="语气: objective | analytical | opinionated | casual"
    )
    session_id: str | None = Field(None, description="会话 ID (thread_id), 不传则自动生成")
    uploaded_files: list[str] | None = Field(
        None, description="已上传文件 ID 列表 (来自 POST /v1/files), 作为研究数据源"
    )
    multi_agent: bool = Field(
        False,
        description="是否启用多 Agent Supervisor 模式 (默认单图流水线)",
    )
    agent_role: str | None = Field(
        None,
        description=(
            "用户可注入行业 persona 字符串, "
            "优先级高于 LLM 动态生成 (AgentCreator). "
            "行业适配采用 4 层机制, 不再使用行业分类器."
        ),
    )
    language: str | None = Field(
        None,
        description=(
            "报告语言代码, 可选 zh (中文, 默认) | en (英文). "
            "不传或传 None 时降级为 settings.report_language. "
            "影响 ReportGenerator 子主题/引言/章节/结论的撰写语言."
        ),
    )
    query_domains: list[str] | None = Field(
        None,
        description="域名过滤白名单, 仅检索这些域名的结果",
    )
    # SELF_HOST=False 时点数校验/扣除所需参数
    # 优先级: org_id > project_id (二者至少传一个才会触发校验/扣除)
    org_id: str | None = Field(
        None,
        description="组织 ID (用于点数校验, 优先于 project_id, SELF_HOST=False 时启用)",
    )
    project_id: str | None = Field(
        None,
        description="项目 ID (org_id 为空时使用, SELF_HOST=False 时启用)",
    )


class ChatCompletionResponse(BaseModel):
    """OpenAI 兼容非流式响应 (增加 sources 结构化字段)."""

    id: str
    object: str = "chat.completion"
    created: int
    model: str
    choices: list[dict[str, Any]]
    # usage 含 cost_usd (float), 放宽为 dict[str, Any] 兼容真实成本
    usage: dict[str, Any]
    # 显式返回 sources 结构化列表
    # 含 title/url/snippet/score 字段, 测试页面与下游消费者可直接渲染
    sources: list[dict[str, Any]] = []
    # 报告输出格式 (markdown/html/pdf/docx/json), 客户端据此选择渲染/下载策略
    report_format: str | None = None
    # PDF 等二进制格式生成的文件路径, 客户端可通过 /v1/reports/{report_id}/download 获取
    file_path: str | None = None
    # 报告主键 UUID (一个 session 可生成多个报告, 下载基于 report_id)
    # 客户端可用此值构造 /v1/reports/{report_id}/download 链接获取多格式输出
    report_id: str | None = None


# ========== 会话持久化辅助 (以 UserId 为单位的对话保存) ==========


async def _persist_user_message(
    session_id: str,
    agent_id: str,
    user_id: str,
    query: str,
) -> None:
    """保存用户消息到 chat_messages, 并确保 research_sessions 记录存在.

    在 chat_completions 端点处理请求前调用:
    - ensure_session: 不存在则创建 research_sessions (含 query + title + client_ip)
    - save_message: 保存 user 消息到 chat_messages

    失败仅告警, 不阻断主流程 (消息持久化为辅助功能).
    """
    try:
        from src.api.middleware import get_request_client_ip
        from src.memory.session_store import get_session_store

        store = get_session_store()
        # 获取客户端 IP (审计追溯用, 从 contextvars 恢复)
        client_ip = get_request_client_ip()
        # 确保会话记录存在 (首次对话时创建, 已存在则更新 query + updated_at)
        title = query[:100] if query else ""
        await store.ensure_session(session_id, agent_id, user_id, query=query, client_ip=client_ip)
        # 若标题为空, 更新为 query 前 100 字符
        if title:
            existing_title = await store.get_session_title(session_id, agent_id, user_id)
            if not existing_title:
                await store.update_session_title(session_id, agent_id, user_id, title)
        # 保存 user 消息
        await store.save_message(
            session_id=session_id,
            agent_id=agent_id,
            user_id=user_id,
            role="user",
            content=query,
        )
    except Exception:
        logger.warning(
            "保存用户消息失败 (不阻断主流程, session=%s): ",
            session_id,
            exc_info=True,
        )


async def _persist_assistant_message(
    session_id: str,
    agent_id: str,
    user_id: str,
    content: str,
    metadata: dict[str, Any] | None = None,
) -> None:
    """保存 assistant 消息到 chat_messages.

    在流式/非流式响应完成后调用 (后台任务, 不阻塞用户响应).
    失败仅告警, 不阻断主流程.

    Args:
        session_id: 会话 ID
        agent_id: Agent 名称
        user_id: 用户 ID
        content: assistant 响应内容
        metadata: 可选元数据 (如 sources, report_id)
    """
    try:
        from src.memory.session_store import get_session_store

        store = get_session_store()
        await store.save_message(
            session_id=session_id,
            agent_id=agent_id,
            user_id=user_id,
            role="assistant",
            content=content,
            metadata=metadata,
        )
    except Exception:
        logger.warning(
            "保存 assistant 消息失败 (不阻断主流程, session=%s): ",
            session_id,
            exc_info=True,
        )


# ========== 端点实现 ==========


@router.post("/chat/completions")
async def chat_completions(
    request: ChatCompletionRequest,
    req: Request,
    authorization: str | None = Header(None),
) -> Any:
    """OpenAI 兼容研究端点.

    测试页面统一调用此端点, 请求体带 stream: true.
    """
    settings = get_settings()

    # 提取最后一条 user 消息作为研究查询
    user_messages = [m for m in request.messages if m.role == "user"]
    if not user_messages:
        raise HTTPException(status_code=400, detail="messages 必须包含至少一条 user 消息")

    query = user_messages[-1].content
    if not query.strip():
        raise HTTPException(status_code=400, detail="查询内容不能为空")

    # 会话 ID (thread_id): 优先请求体, 其次中间件, 最后生成
    session_id = request.session_id or get_request_session_id() or str(uuid.uuid4())

    # user_id / agent_id 由中间件注入
    user_id = get_request_user_id()
    agent_id = get_request_agent_id()

    # 会话持久化: 保存 user 消息到 chat_messages + 确保 research_sessions 存在
    # (以 UserId 为单位的会话持久化, 失败不阻断主流程)
    if user_id and agent_id:
        await _persist_user_message(session_id, agent_id, user_id, query)

    # 查询意图分类 (短查询 + 离题闲聊 + 对话/研究路由)
    # 先检查会话是否已有报告 (用于分类器 has_report 参数)
    has_report = await _has_report(session_id)
    intent = await get_query_intent_classifier().classify(query, has_report)

    # CHAT 首轮保护 — 无已有报告时降级 OFF_TOPIC
    # 避免首轮闲聊走 chat graph 消耗 SMART LLM; 显式 report_type 时强制走 researcher graph
    if (
        intent == QueryIntent.CHAT
        and request.report_type is None
        and settings.chat_requires_report
        and not has_report
    ):
        logger.debug(
            "CHAT 首轮无报告, 降级 OFF_TOPIC (session_id=%s)",
            session_id,
        )
        intent = QueryIntent.OFF_TOPIC

    # 短查询/离题闲聊 — 走 ChitchatResponder
    # CHITCHAT_FAST_LLM_OPTIMIZATION_PLAN.md: 始终走 FAST_LLM 人性化响应, FAST 失败时降级 multi-template
    if intent in (QueryIntent.SHORT_QUERY, QueryIntent.OFF_TOPIC):
        # ChitchatResponder (FAST_LLM + Persona + 三段式 + 多模板兜底)
        category = (
            _infer_off_topic_category(query) if intent == QueryIntent.OFF_TOPIC else "greeting"
        )
        responder = get_chitchat_responder()
        # 会话持久化: 从 checkpointer 加载历史并注入 ChitchatResponder
        history = await _load_chitchat_history(session_id)
        if intent == QueryIntent.SHORT_QUERY:
            if request.stream:
                return StreamingResponse(
                    _stream_chitchat(
                        responder.respond_short_query(
                            query,
                            session_id=session_id,
                            user_id=user_id,
                            stream=True,
                            history=history,
                        ),
                        request,
                        session_id,
                        intent=intent.value,
                        save_query=query,
                        save_user_id=user_id,
                        save_agent_id=agent_id,
                    ),
                    media_type="text/event-stream",
                    headers={
                        "Cache-Control": "no-cache",
                        "Connection": "keep-alive",
                        "X-Session-Id": session_id,
                    },
                )
            response = await _run_chitchat(
                responder.respond_short_query(
                    query,
                    session_id=session_id,
                    user_id=user_id,
                    stream=False,
                    history=history,
                ),
                request,
                session_id,
                intent=intent.value,
            )
            # 会话持久化: 保存 ChitchatResponder 响应回 checkpointer + chat_messages
            chitchat_content = response.choices[0]["message"]["content"]
            await _save_chitchat_response(
                session_id,
                query,
                chitchat_content,
                user_id=user_id,
                agent_id=agent_id,
            )
            if user_id and agent_id:
                await _persist_assistant_message(
                    session_id,
                    agent_id,
                    user_id,
                    chitchat_content,
                    metadata={"intent": intent.value},
                )
            return _with_session_id(response, session_id)
        else:  # OFF_TOPIC
            if request.stream:
                return StreamingResponse(
                    _stream_chitchat(
                        responder.respond_off_topic(
                            query,
                            category=category,
                            session_id=session_id,
                            user_id=user_id,
                            stream=True,
                            history=history,
                        ),
                        request,
                        session_id,
                        intent=intent.value,
                        save_query=query,
                        save_user_id=user_id,
                        save_agent_id=agent_id,
                    ),
                    media_type="text/event-stream",
                    headers={
                        "Cache-Control": "no-cache",
                        "Connection": "keep-alive",
                        "X-Session-Id": session_id,
                    },
                )
            response = await _run_chitchat(
                responder.respond_off_topic(
                    query,
                    category=category,
                    session_id=session_id,
                    user_id=user_id,
                    stream=False,
                    history=history,
                ),
                request,
                session_id,
                intent=intent.value,
            )
            # 会话持久化: 保存 ChitchatResponder 响应回 checkpointer + chat_messages
            chitchat_content = response.choices[0]["message"]["content"]
            await _save_chitchat_response(
                session_id,
                query,
                chitchat_content,
                user_id=user_id,
                agent_id=agent_id,
            )
            if user_id and agent_id:
                await _persist_assistant_message(
                    session_id,
                    agent_id,
                    user_id,
                    chitchat_content,
                    metadata={"intent": intent.value},
                )
            return _with_session_id(response, session_id)

    # CHAT 意图走 chat graph (仅追问场景, 首轮已被降级到 OFF_TOPIC)
    # 显式指定 report_type 时强制走 researcher graph (用户明确要新研究)
    if intent == QueryIntent.CHAT and request.report_type is None:
        # 走 chat graph (复用会话历史 + report_md 上下文, 首轮时 report_md 为空)
        # 注意: initial_state 不含 report_md (由 checkpointer 自动加载, 避免覆盖)
        chat_state: dict[str, Any] = {
            "query": query,
            "session_id": session_id,
            "user_id": user_id,
            "agent_id": agent_id,
            "query_intent": intent.value,
            "messages": [],
        }
        chat_config: dict[str, Any] = {"configurable": {"thread_id": session_id}}
        if request.stream:
            return StreamingResponse(
                _stream_chat(
                    chat_state,
                    chat_config,
                    request,
                    session_id,
                    save_user_id=user_id,
                    save_agent_id=agent_id,
                ),
                media_type="text/event-stream",
                headers={
                    "Cache-Control": "no-cache",
                    "Connection": "keep-alive",
                    "X-Session-Id": session_id,
                },
            )
        else:
            chat_response = await _run_chat(chat_state, chat_config, request, session_id)
            # 会话持久化: 保存 assistant 消息到 chat_messages
            if user_id and agent_id:
                chat_content = chat_response.choices[0]["message"]["content"]
                if chat_content:
                    await _persist_assistant_message(
                        session_id,
                        agent_id,
                        user_id,
                        chat_content,
                        metadata={"intent": "chat"},
                    )
            return _with_session_id(chat_response, session_id)

    # RESEARCH 意图 (或 CHAT + 显式 report_type) → researcher graph
    # SELF_HOST=False 时, 进入研究前校验 Agent 点数
    # token 不得入日志/持久化; 仅 RESEARCH 意图校验/扣除
    if not settings.self_host and (request.org_id or request.project_id):
        from src.api.agentinsight_client import get_agentinsight_client

        # 提取 token (从 Authorization Bearer 头)
        token = ""
        if authorization and authorization.lower().startswith("bearer "):
            token = authorization[7:].strip()

        if not token:
            raise HTTPException(
                status_code=401,
                detail="缺少 Authorization Bearer Token (SELF_HOST=False 模式必需)",
            )

        client = get_agentinsight_client()
        exceeded, _err = await client.validate_agent_usage(token)
        if exceeded:
            raise HTTPException(
                status_code=429,
                detail="本月 Agent 调用次数已达上限, 请联系管理员升级套餐",
            )

    # 报告配置 (新研究模式)
    report_type = request.report_type or settings.default_report_type
    report_format = request.report_format or settings.default_report_format
    tone = request.tone or "objective"
    # report_type == "deep_research" 时启用递归深度研究
    # summary 和 subtopics 正确映射到对应 research_mode, 避免被降级为 basic
    if report_type == "deep_research":
        research_mode = "deep"
    elif report_type in ("summary", "subtopics"):
        research_mode = report_type  # 让 ResearchConductor 走 _conduct_summary / _conduct_subtopics
    else:
        research_mode = "basic"

    # 加载已上传文件上下文 (用户需求 8)
    uploaded_files_context: list[str] = []
    if request.uploaded_files:
        # _load_uploaded_files_context 为 async, 内部文件 I/O 经 asyncio.to_thread
        uploaded_files_context = await _load_uploaded_files_context(
            request.uploaded_files, user_id, agent_id
        )

    # 初始 State (TypedDict, 节点返回 delta)
    initial_state: dict[str, Any] = {
        "query": query,
        "session_id": session_id,
        "user_id": user_id,
        "agent_id": agent_id,
        "query_intent": intent.value,
        "report_type": report_type,
        "report_format": report_format,
        "tone": tone,
        "uploaded_files_context": uploaded_files_context,
        "messages": [],
        # agent_role (来自 ChatRequest 或 settings) 优先级高于 LLM 动态生成
        "agent_role": request.agent_role or settings.agent_role or "",
        "agent_role_server": "",
        # 报告语言 (来自 ChatRequest.language 或 settings.report_language)
        # 可选值: zh (中文, 默认) | en (英文)
        "report_language": (request.language or settings.report_language or "zh").lower(),
        # 域名过滤白名单
        "query_domains": request.query_domains or [],
        "sub_queries": [],
        "contexts": [],
        "sources": [],
        "visited_urls": [],
        "curated_sources": [],
        "report_md": "",  # 兼容字段, 新代码用 report_formats
        "report_formats": {},  # {md|html|pdf|docx|json: 内容或路径}
        "status": "pending",
        # 深度研究配置
        "research_mode": research_mode,
        "deep_research_breadth": settings.deep_research_breadth,
        "deep_research_depth": settings.deep_research_depth,
        # 多 Agent 迭代计数器 (Annotated[int, operator.add] 累加)
        "iteration_count": 0,
    }

    # LangGraph 配置: thread_id 做会话隔离
    graph_config = {"configurable": {"thread_id": session_id}}

    # IP-based 用户每日报告限额检查 (仅 self_host=True + IP-based 用户)
    # 限额从数据库 report_limits 表读取 (已从环境变量迁移)
    if settings.self_host and user_id.startswith("ip_"):
        from src.api.ip_user_resolver import check_daily_report_limit

        allowed, current_count, effective_limit = await check_daily_report_limit(user_id, agent_id)
        if not allowed:
            limit_msg = (
                f"您今日的报告生成次数已达上限 ({current_count}/{effective_limit})。"
                f"每日限额将在北京时间次日 0 点重置, 届时可继续使用。"
            )
            if request.stream:
                # 流式: 发送友好提示后关闭 SSE
                async def _limit_exceeded_stream() -> Any:
                    completion_id = f"chatcmpl-{uuid.uuid4().hex[:24]}"
                    chunk = {
                        "id": completion_id,
                        "object": "chat.completion.chunk",
                        "created": int(time.time()),
                        "model": request.model,
                        "choices": [
                            {
                                "index": 0,
                                "delta": {"content": limit_msg},
                                "finish_reason": None,
                            }
                        ],
                    }
                    yield f"data: {orjson.dumps(chunk).decode('utf-8')}\n\n"
                    done_chunk = {
                        "id": completion_id,
                        "object": "chat.completion.chunk",
                        "created": int(time.time()),
                        "model": request.model,
                        "choices": [
                            {
                                "index": 0,
                                "delta": {},
                                "finish_reason": "stop",
                            }
                        ],
                    }
                    yield f"data: {orjson.dumps(done_chunk).decode('utf-8')}\n\n"
                    yield "data: [DONE]\n\n"

                return StreamingResponse(
                    _limit_exceeded_stream(),
                    media_type="text/event-stream",
                    headers={
                        "Cache-Control": "no-cache",
                        "Connection": "keep-alive",
                        "X-Session-Id": session_id,
                    },
                )
            else:
                # 非流式: 返回 JSON 错误
                return JSONResponse(
                    status_code=429,
                    content={
                        "error": {
                            "message": limit_msg,
                            "type": "rate_limit_exceeded",
                            "code": "daily_report_limit",
                        }
                    },
                    headers={"X-Session-Id": session_id},
                )

    if request.stream:
        return StreamingResponse(
            _stream_research(
                initial_state,
                graph_config,
                request,
                session_id,
                authorization,
                save_user_id=user_id,
                save_agent_id=agent_id,
            ),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "X-Session-Id": session_id,
            },
        )
    else:
        # 非流式: 完整执行后返回
        research_response = await _run_research(
            initial_state,
            graph_config,
            request,
            session_id,
            authorization,
        )
        # 会话持久化: 保存 assistant 消息到 chat_messages
        if user_id and agent_id:
            research_content = research_response.choices[0]["message"]["content"]
            if research_content:
                await _persist_assistant_message(
                    session_id,
                    agent_id,
                    user_id,
                    research_content,
                    metadata={
                        "intent": "research",
                        "report_id": research_response.report_id,
                        "sources": research_response.sources[:10]
                        if research_response.sources
                        else [],
                    },
                )
        return _with_session_id(research_response, session_id)


def _with_session_id(response: ChatCompletionResponse, session_id: str) -> JSONResponse:
    """将非流式 ChatCompletionResponse 包装为 JSONResponse 并注入 X-Session-Id 响应头.

    流式分支 (StreamingResponse) 已在 headers 中显式设置 X-Session-Id;
    非流式分支返回 Pydantic 模型时 FastAPI 自动转 JSONResponse 会丢失自定义头,
    本函数统一在非流式路径注入 X-Session-Id (会话持久化回归测试依赖此头).

    Args:
        response: ChatCompletionResponse Pydantic 模型
        session_id: 会话 ID

    Returns:
        JSONResponse 含 X-Session-Id 响应头
    """
    return JSONResponse(
        content=response.model_dump(),
        headers={"X-Session-Id": session_id},
    )


async def _stream_research(
    initial_state: dict[str, Any],
    graph_config: dict[str, Any],
    request: ChatCompletionRequest,
    session_id: str,
    authorization: str | None = None,
    *,
    save_user_id: str | None = None,
    save_agent_id: str | None = None,
) -> Any:
    """流式 SSE 响应生成器.

    接入 LangGraph astream, 逐节点 yield 进度 + 最终报告.
    用 trace_agent 包裹 graph.ainvoke 作为根 span.

    会话持久化: 流式结束后保存完整报告内容到 chat_messages
    (save_user_id + save_agent_id 非空时触发).
    """
    completion_id = f"chatcmpl-{uuid.uuid4().hex[:24]}"
    created = int(time.time())

    def _sse_chunk(delta: dict[str, Any], finish_reason: str | None = None) -> str:
        """构造 SSE 数据帧."""
        chunk = {
            "id": completion_id,
            "object": "chat.completion.chunk",
            "created": created,
            "model": request.model,
            "choices": [
                {
                    "index": 0,
                    "delta": delta,
                    "finish_reason": finish_reason,
                }
            ],
        }
        return f"data: {orjson.dumps(chunk).decode('utf-8')}\n\n"

    # SSE 首块 (role)
    yield _sse_chunk({"role": "assistant"})

    # 收集流式内容 (用于会话持久化)
    collected_content: list[str] = []

    # 根 span 包裹整次研究
    async with trace_agent(
        name="agentinsight-researcher",
        input={"query": initial_state["query"][:200], "session_id": session_id},
        metadata={
            "session_id": session_id,
            "intent": "research",
            "user_id": initial_state.get("user_id"),
        },
        session_id=session_id,
        user_id=initial_state.get("user_id"),
    ):
        try:
            graph = await _get_graph(multi_agent=request.multi_agent)

            # 节点名称中文映射 (用于流式进度提示)
            node_label = {
                "agent_creator": "生成研究角色",
                "deep_research": "深度研究",
                "research_conductor": "并行检索研究",
                "source_curator": "来源策展",
                "report_generator": "生成报告",
                "publisher": "格式化输出",
                # 多 Agent 节点
                "supervisor": "调度决策",
                "reviewer": "评估来源",
                "writer": "撰写报告",
            }

            final_state: dict[str, Any] = {}

            # astream 流式输出节点进度
            async for event in graph.astream(
                initial_state, config=graph_config, stream_mode="updates"
            ):
                # event 格式: {node_name: delta_dict}
                for node_name, delta in event.items():
                    if not isinstance(delta, dict):
                        continue
                    final_state.update(delta)

                    # 节点开始进度提示
                    label = node_label.get(node_name, node_name)
                    progress = f"\n\n> **[{label}]** "
                    if node_name == "agent_creator" and delta.get("agent_role_server"):
                        progress += f"已生成研究角色: {delta['agent_role_server']}\n"
                    elif node_name == "research_conductor":
                        sq_count = len(delta.get("sub_queries", []))
                        ctx_count = len(delta.get("contexts", []))
                        src_count = len(delta.get("sources", []))
                        progress += (
                            f"已生成 {sq_count} 子查询, 采集 {ctx_count} 上下文, {src_count} 来源\n"
                        )
                    elif node_name == "source_curator":
                        if delta.get("curated_sources"):
                            progress += f"已策展 {len(delta['curated_sources'])} 来源\n"
                    elif node_name == "report_generator" and delta.get("report_md"):
                        # SELF_HOST=False 时, 报告生成成功后异步扣除点数 (不阻塞流式响应)
                        # token 不得入日志/持久化; 仅 RESEARCH 意图扣除
                        settings_priv = get_settings()
                        if not settings_priv.self_host and (request.org_id or request.project_id):
                            from src.api.agentinsight_client import get_agentinsight_client

                            token = (
                                authorization[7:].strip()
                                if authorization and authorization.lower().startswith("bearer ")
                                else ""
                            )
                            if token:
                                # 保留 Task 引用防止 GC 取消 (扣点任务有副作用)
                                _create_background_task(
                                    get_agentinsight_client().deduct_agent_usage(token)
                                )
                        # 非 markdown 格式: 跳过逐段推送 (publisher 节点会推送最终格式)
                        fmt = request.report_format or "markdown"
                        if fmt != "markdown":
                            # 仅推送进度提示, 不推送 markdown 内容
                            progress += f"报告生成完成 ({fmt} 格式转换中...)\n"
                            yield _sse_chunk({"progress": progress})
                            progress = ""
                            continue
                        # markdown 格式: 维持逐段流式推送
                        report_md = delta["report_md"]
                        # 分块输出报告 (按段落)
                        paragraphs = report_md.split("\n\n")
                        for para in paragraphs:
                            chunk_text = para + "\n\n"
                            collected_content.append(chunk_text)
                            yield _sse_chunk({"content": chunk_text})
                            # 背压: 让出控制权, 允许消费者 (SSE writer) 刷出缓冲
                            await asyncio.sleep(0.001)
                        continue  # 跳过下面的进度提示
                    elif node_name == "publisher":
                        fmt = delta.get("report_format", "markdown")
                        progress += f"已发布 ({fmt})\n"
                        # 推送 report_id (前端用于构造 /v1/reports/{report_id}/download 链接)
                        if delta.get("report_id"):
                            yield _sse_chunk({"report_id": delta["report_id"]})
                        # 流式推送最终格式的报告内容 (统一从 report_formats 读取)
                        new_formats = delta.get("report_formats") or {}
                        if new_formats.get("html"):
                            collected_content.append(new_formats["html"])
                            yield _sse_chunk({"content": new_formats["html"]})
                        elif new_formats.get("json"):
                            collected_content.append(new_formats["json"])
                            yield _sse_chunk({"content": new_formats["json"]})
                        elif new_formats.get("pdf"):
                            yield _sse_chunk(
                                {"file_path": new_formats["pdf"], "report_format": "pdf"}
                            )
                        # docx 二进制不适合 SSE, 跳过 (客户端可走非流式或下载端点)
                        continue

                    yield _sse_chunk({"content": progress})
                    if progress:
                        collected_content.append(progress)
                    # 背压: 让出控制权, 允许消费者 (SSE writer) 刷出缓冲
                    await asyncio.sleep(0.001)

            # 输出最终报告元信息
            fmt = request.report_format or "markdown"
            # 优先从 report_formats 读取, 兼容期回退旧字段
            final_formats = final_state.get("report_formats") or {}
            final_content = (
                final_formats.get("html")
                or final_formats.get("json")
                or final_formats.get("md")
                or final_state.get("report_md")
                or ""
            )
            if not final_content:
                _empty_hint = "\n\n*未生成报告内容 (可能上下文为空)*"
                collected_content.append(_empty_hint)
                yield _sse_chunk({"content": _empty_hint})
            # 如果是 PDF, 推送 file_path
            pdf_path = final_formats.get("pdf") or final_state.get("report_pdf_path")
            if fmt == "pdf" and pdf_path:
                yield _sse_chunk({"file_path": pdf_path, "report_format": "pdf"})

            # 流式推送 sources 元信息帧 (在 finish 之前)
            # 客户端可通过 delta.sources 获取结构化来源列表
            final_sources = final_state.get("curated_sources") or final_state.get("sources", [])
            normalized_sources: list[dict[str, Any]] = []
            for src in final_sources[:20]:
                if isinstance(src, dict):
                    normalized_sources.append(
                        {
                            "title": src.get("title", "") or "",
                            "url": src.get("url", "") or src.get("href", "") or "",
                            "snippet": src.get("snippet", "") or "",
                            "score": src.get("score", 0.0) or 0.0,
                        }
                    )
            if normalized_sources:
                yield _sse_chunk({"sources": normalized_sources})

            # 报告持久化移到后台任务 (不阻塞 [DONE], 先 yield [DONE] 再后台持久化)
            # 保留 Task 引用防止 GC 取消
            async def _persist_report() -> None:
                try:
                    from src.memory.report_store import get_report_store

                    _report_store = get_report_store()
                    _saved_report_id = await _report_store.save_report(
                        session_id=session_id,
                        user_id=initial_state.get("user_id", ""),
                        agent_id=initial_state.get("agent_id", ""),
                        query=initial_state.get("query", ""),
                        report_md=final_state.get("report_md", ""),
                        report_format=final_state.get("report_format", "markdown"),
                        sources=(
                            final_state.get("curated_sources") or final_state.get("sources", [])
                        ),
                        agent_role=final_state.get("agent_role_server"),
                    )
                    if _saved_report_id:
                        # IP-based 用户报告生成成功后异步递增每日计数
                        _uid = initial_state.get("user_id", "")
                        _aid = initial_state.get("agent_id", "")
                        if _uid.startswith("ip_") and _aid:
                            from src.api.ip_user_resolver import increment_daily_report_count

                            await increment_daily_report_count(_uid, _aid)
                except Exception:
                    logger.warning("报告持久化存储失败 (不阻断主流程)", exc_info=True)

            _create_background_task(_persist_report())

        except Exception as e:
            logger.exception("研究流水线执行失败")
            _err_msg = f"\n\n**研究执行失败**: {str(e)[:200]}"
            collected_content.append(_err_msg)
            yield _sse_chunk({"content": _err_msg})
            # 会话持久化: 保存 assistant 消息到 chat_messages (含错误消息)
            # 必须在 yield [DONE] 之前, 否则 SSE 连接中断后 yield 之后的代码不执行
            if save_user_id and save_agent_id and collected_content:
                full_response = "".join(collected_content)
                if full_response:
                    await _persist_assistant_message(
                        session_id,
                        save_agent_id,
                        save_user_id,
                        full_response,
                        metadata={"intent": "research", "status": "error"},
                    )
            # 错误时 finish_reason="error" (客户端可据此区分失败)
            yield _sse_chunk({}, finish_reason="error")
            yield "data: [DONE]\n\n"
            return

    # 会话持久化: 保存 assistant 消息到 chat_messages (必须在 yield [DONE] 之前,
    # 否则 SSE 连接中断后 yield 之后的代码不执行, 导致历史消息丢失)
    if save_user_id and save_agent_id and collected_content:
        full_response = "".join(collected_content)
        if full_response:
            await _persist_assistant_message(
                session_id,
                save_agent_id,
                save_user_id,
                full_response,
                metadata={"intent": "research"},
            )

    # SSE 末块 (finish_reason)
    yield _sse_chunk({}, finish_reason="stop")
    yield "data: [DONE]\n\n"


async def _run_research(
    initial_state: dict[str, Any],
    graph_config: dict[str, Any],
    request: ChatCompletionRequest,
    session_id: str,
    authorization: str | None = None,
) -> ChatCompletionResponse:
    """非流式研究执行.

    接入 LangGraph ainvoke.
    trace_agent 根 span.
    """
    completion_id = f"chatcmpl-{uuid.uuid4().hex[:24]}"
    created = int(time.time())

    async with trace_agent(
        name="agentinsight-researcher",
        input={"query": initial_state["query"][:200], "session_id": session_id},
        metadata={
            "session_id": session_id,
            "intent": "research",
            "user_id": initial_state.get("user_id"),
        },
        session_id=session_id,
        user_id=initial_state.get("user_id"),
    ):
        file_path: str | None = None
        report_id: str | None = None
        settings = get_settings()
        try:
            graph = await _get_graph(multi_agent=request.multi_agent)
            # graph.ainvoke 超时保护 (防止节点卡死永久挂起)
            final_state = await asyncio.wait_for(
                graph.ainvoke(initial_state, config=graph_config),
                timeout=settings.graph_total_timeout,
            )
            fmt = request.report_format or "markdown"
            # 优先从 report_formats 读取, 兼容期回退旧字段
            final_formats = final_state.get("report_formats") or {}
            if fmt == "html" and final_formats.get("html"):
                content = final_formats["html"]
            elif fmt == "json" and final_formats.get("json"):
                content = final_formats["json"]
            elif fmt == "pdf" and final_formats.get("pdf"):
                # PDF 返回路径信息, 客户端可通过下载端点获取
                content = final_formats["pdf"]
                file_path = final_formats["pdf"]
            elif fmt == "docx" and final_formats.get("docx"):
                # DOCX 二进制无法直接放入 OpenAI 兼容响应, 返回提示信息
                content = "[DOCX 报告已生成, 请通过下载端点获取]"
            else:
                content = final_formats.get("md") or final_state.get("report_md", "")
            if not content:
                content = "未生成报告内容 (可能上下文为空)"
            # 报告持久化 (从 publisher_node 移到 API 层, 节点纯函数无副作用)
            # graph 完成后调用 report_store.save_report, 保存失败仅 warn 不影响响应
            # (用户已收到报告内容, 不返回 500; 节点纯函数约束)
            # 客户端用 report_id 构造 /v1/reports/{report_id}/download 链接
            # mypy no-redef: 第 706 行已声明 report_id, 此处重新赋值 (非重新声明)
            report_id = None
            try:
                from src.memory.report_store import get_report_store

                _report_store = get_report_store()
                _saved = await _report_store.save_report(
                    session_id=session_id,
                    user_id=initial_state.get("user_id", ""),
                    agent_id=initial_state.get("agent_id", ""),
                    query=initial_state.get("query", ""),
                    report_md=final_state.get("report_md", ""),
                    report_format=final_state.get("report_format", "markdown"),
                    sources=(final_state.get("curated_sources") or final_state.get("sources", [])),
                    agent_role=final_state.get("agent_role_server"),
                )
                if _saved:
                    report_id = _saved
            except Exception:
                logger.warning("报告持久化存储失败 (不阻断主流程)", exc_info=True)

            # sources 作为结构化字段返回 (优先 curated_sources, 回退 sources)
            sources = final_state.get("curated_sources") or final_state.get("sources", [])
            # 规范化: 仅保留 title/url/snippet/score 四字段, 便于客户端渲染
            normalized_sources: list[dict[str, Any]] = []
            for src in sources[:20]:
                if isinstance(src, dict):
                    normalized_sources.append(
                        {
                            "title": src.get("title", "") or "",
                            "url": src.get("url", "") or src.get("href", "") or "",
                            "snippet": src.get("snippet", "") or "",
                            "score": src.get("score", 0.0) or 0.0,
                        }
                    )

            # 读取 LLMClient 真实成本 (优先于字符估算)
            total_cost_usd = final_state.get("total_cost_usd", 0.0) or 0.0
            total_tokens = final_state.get("total_tokens", 0) or 0
            token_logs = final_state.get("token_logs", []) or []
            if total_tokens > 0:
                prompt_tokens = sum(
                    int(log.get("prompt_tokens", 0)) for log in token_logs if isinstance(log, dict)
                )
                completion_tokens = sum(
                    int(log.get("completion_tokens", 0))
                    for log in token_logs
                    if isinstance(log, dict)
                )
                usage = {
                    "prompt_tokens": prompt_tokens,
                    "completion_tokens": completion_tokens,
                    "total_tokens": total_tokens,
                    "cost_usd": round(total_cost_usd, 6),
                }
            else:
                # 降级: 字符估算 (向后兼容)
                prompt_tokens = len(initial_state["query"]) // 4
                completion_tokens = len(content) // 4
                usage = {
                    "prompt_tokens": prompt_tokens,
                    "completion_tokens": completion_tokens,
                    "total_tokens": prompt_tokens + completion_tokens,
                }

        except Exception as e:
            logger.exception("研究流水线执行失败")
            content = f"研究执行失败: {str(e)[:500]}"
            normalized_sources = []
            usage = {
                "prompt_tokens": len(initial_state["query"]) // 4,
                "completion_tokens": len(content) // 4,
                "total_tokens": (len(initial_state["query"]) + len(content)) // 4,
            }

    # SELF_HOST=False 时, 报告生成成功后同步扣除点数 (不阻断响应)
    # token 不得入日志/持久化; 仅 RESEARCH 意图扣除
    settings_priv = get_settings()
    if (
        not settings_priv.self_host
        and (request.org_id or request.project_id)
        and content
        and "研究执行失败" not in content
    ):
        try:
            from src.api.agentinsight_client import get_agentinsight_client

            token = (
                authorization[7:].strip()
                if authorization and authorization.lower().startswith("bearer ")
                else ""
            )
            if token:
                await get_agentinsight_client().deduct_agent_usage(token)
        except Exception as e:  # noqa: BLE001
            logger.warning("扣除点数失败 (不阻断响应): %s", e)

    # IP-based 用户报告生成成功后递增每日计数 (self_host=True 且 user_id 以 "ip_" 前缀)
    # 仅报告内容非空且未失败时计数; 失败/异常不计数 (符合"需生成报告成功才计数"需求)
    _uid_final = initial_state.get("user_id", "")
    _aid_final = initial_state.get("agent_id", "")
    if content and "研究执行失败" not in content and _uid_final.startswith("ip_") and _aid_final:
        try:
            from src.api.ip_user_resolver import increment_daily_report_count

            await increment_daily_report_count(_uid_final, _aid_final)
        except Exception as e:  # noqa: BLE001
            logger.warning("每日报告计数递增失败 (不阻断响应): %s", e)

    return ChatCompletionResponse(
        id=completion_id,
        created=created,
        model=request.model,
        choices=[
            {
                "index": 0,
                "message": {"role": "assistant", "content": content},
                "finish_reason": "stop",
            }
        ],
        usage=usage,
        sources=normalized_sources,
        report_format=request.report_format or "markdown",
        file_path=file_path,
        report_id=report_id,
    )


# ========== 短查询保护 (不走任何 graph) + 离题闲聊保护 ==========


async def _stream_short_query(
    reply: str,
    request: ChatCompletionRequest,
    session_id: str,
    *,
    intent: str = "short_query",
) -> Any:
    """流式 SSE 短查询/离题回复生成器.

    直接返回 settings.short_query_reply / settings.off_topic_reply, 不走任何 graph.
    用 trace_agent 包裹作为根 span.

    Args:
        reply: 回复内容 (短查询回复语或离题回复语)
        request: OpenAI 兼容请求
        session_id: 会话 ID
        intent: 意图类型 ("short_query" 或 "off_topic"), 用于 trace 区分
    """
    completion_id = f"chatcmpl-{uuid.uuid4().hex[:24]}"
    created = int(time.time())

    def _sse_chunk(delta: dict[str, Any], finish_reason: str | None = None) -> str:
        """构造 SSE 数据帧."""
        chunk = {
            "id": completion_id,
            "object": "chat.completion.chunk",
            "created": created,
            "model": request.model,
            "choices": [
                {
                    "index": 0,
                    "delta": delta,
                    "finish_reason": finish_reason,
                }
            ],
        }
        return f"data: {orjson.dumps(chunk).decode('utf-8')}\n\n"

    # SSE 首块 (role)
    yield _sse_chunk({"role": "assistant"})

    span_name = f"agentinsight-researcher-{intent}"
    async with trace_agent(
        name=span_name,
        input={"session_id": session_id},
        metadata={
            "session_id": session_id,
            "intent": intent,
        },
        session_id=session_id,
    ):
        yield _sse_chunk({"content": reply})

    # SSE 末块 (finish_reason)
    yield _sse_chunk({}, finish_reason="stop")
    yield "data: [DONE]\n\n"


async def _run_short_query(
    reply: str,
    request: ChatCompletionRequest,
    session_id: str,
    *,
    intent: str = "short_query",
) -> ChatCompletionResponse:
    """非流式短查询/离题回复.

    直接返回 settings.short_query_reply / settings.off_topic_reply, 不走任何 graph.
    trace_agent 根 span.

    Args:
        reply: 回复内容 (短查询回复语或离题回复语)
        request: OpenAI 兼容请求
        session_id: 会话 ID
        intent: 意图类型 ("short_query" 或 "off_topic"), 用于 trace 区分
    """
    completion_id = f"chatcmpl-{uuid.uuid4().hex[:24]}"
    created = int(time.time())

    span_name = f"agentinsight-researcher-{intent}"
    async with trace_agent(
        name=span_name,
        input={"session_id": session_id},
        metadata={
            "session_id": session_id,
            "intent": intent,
        },
        session_id=session_id,
    ):
        content = reply

    prompt_tokens = 0
    completion_tokens = len(content) // 4

    return ChatCompletionResponse(
        id=completion_id,
        created=created,
        model=request.model,
        choices=[
            {
                "index": 0,
                "message": {"role": "assistant", "content": content},
                "finish_reason": "stop",
            }
        ],
        usage={
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": prompt_tokens + completion_tokens,
        },
    )


# ========== ChitchatResponder 辅助函数 (CHITCHAT_FAST_LLM_OPTIMIZATION_PLAN.md) ==========


def _infer_off_topic_category(query: str) -> str:
    """从用户查询推断 OFF_TOPIC 子类.

    关键词匹配实现. 8 类: greeting/identity/emotion/entertainment/common_sense/capability_check/topic_switch/evaluation.
    """
    q = query.lower().strip()
    # 身份询问
    if any(
        kw in q
        for kw in (
            "你叫什么",
            "你是谁",
            "你多大了",
            "你是真人",
            "你是机器人",
            "什么模型",
            "你的名字",
        )
    ):
        return "identity"
    # 情绪表达
    if any(kw in q for kw in ("我好累", "我好开心", "心情", "难过", "开心", "烦", "生气", "无聊")):
        return "emotion"
    # 娱乐请求
    if any(kw in q for kw in ("讲个笑话", "说个故事", "唱首歌", "玩游戏", "讲笑话")):
        return "entertainment"
    # 常识问题
    if any(kw in q for kw in ("等于几", "天气", "几点", "日期", "今天", "现在时间")):
        return "common_sense"
    # 能力询问
    if any(kw in q for kw in ("你能做什么", "你会什么", "功能", "能力", "怎么用", "帮助")):
        return "capability_check"
    # 话题转移
    if any(kw in q for kw in ("不想", "算了", "退出", "结束", "拜拜", "再见")):
        return "topic_switch"
    # 评价
    if any(kw in q for kw in ("谢谢", "感谢", "很好", "不错", "太棒了", "厉害")):
        return "evaluation"
    # 默认: 问候
    return "greeting"


async def _stream_chitchat(
    content_iterator: Any,
    request: ChatCompletionRequest,
    session_id: str,
    *,
    intent: str = "short_query",
    save_query: str | None = None,
    save_user_id: str | None = None,
    save_agent_id: str | None = None,
) -> Any:
    """流式 SSE 闲聊响应生成器.

    从 ChitchatResponder 的 AsyncIterator 逐块 yield 内容, 包装为 SSE 格式.
    trace_agent 根 span.

    会话持久化:
    - 流式结束后将完整响应保存回 checkpointer (save_query 非空时)
    - 同时保存 assistant 消息到 chat_messages (以 UserId 为单位的会话持久化)

    Args:
        content_iterator: ChitchatResponder 返回的 AsyncIterator[str]
        request: OpenAI 兼容请求
        session_id: 会话 ID
        intent: 意图类型 ("short_query" 或 "off_topic")
        save_query: 用户查询 (非空时触发保存, 会话持久化)
        save_user_id: 用户 ID (保存时注入 state 元数据)
        save_agent_id: Agent ID (保存时注入 state 元数据)
    """
    completion_id = f"chatcmpl-{uuid.uuid4().hex[:24]}"
    created = int(time.time())

    def _sse_chunk(delta: dict[str, Any], finish_reason: str | None = None) -> str:
        """构造 SSE 数据帧."""
        chunk = {
            "id": completion_id,
            "object": "chat.completion.chunk",
            "created": created,
            "model": request.model,
            "choices": [
                {
                    "index": 0,
                    "delta": delta,
                    "finish_reason": finish_reason,
                }
            ],
        }
        return f"data: {orjson.dumps(chunk).decode('utf-8')}\n\n"

    # SSE 首块 (role)
    yield _sse_chunk({"role": "assistant"})

    # 收集流式内容 (用于会话持久化保存)
    collected_content: list[str] = []

    span_name = f"agentinsight-researcher-{intent}"
    async with trace_agent(
        name=span_name,
        input={"session_id": session_id, "intent": intent},
        metadata={
            "session_id": session_id,
            "intent": intent,
            "mode": "chitchat_responder",
        },
        session_id=session_id,
    ):
        # 逐块 yield FAST_LLM 流式内容
        async for chunk in content_iterator:
            if chunk:
                collected_content.append(chunk)
                yield _sse_chunk({"content": chunk})

    # 会话持久化: 流式结束后保存完整响应回 checkpointer + chat_messages
    # 必须在 yield [DONE] 之前, 否则 SSE 连接中断后 yield 之后的代码不执行, 导致历史消息丢失
    if save_query is not None:
        full_response = "".join(collected_content)
        if full_response:
            await _save_chitchat_response(
                session_id,
                save_query,
                full_response,
                user_id=save_user_id,
                agent_id=save_agent_id,
            )
            # 保存 assistant 消息到 chat_messages (以 UserId 为单位的会话持久化)
            if save_user_id and save_agent_id:
                await _persist_assistant_message(
                    session_id,
                    save_agent_id,
                    save_user_id,
                    full_response,
                    metadata={"intent": intent},
                )

    # SSE 末块 (finish_reason)
    yield _sse_chunk({}, finish_reason="stop")
    yield "data: [DONE]\n\n"


async def _run_chitchat(
    content_or_future: Any,
    request: ChatCompletionRequest,
    session_id: str,
    *,
    intent: str = "short_query",
) -> ChatCompletionResponse:
    """非流式闲聊响应.

    从 ChitchatResponder 获取完整字符串响应.
    trace_agent 根 span.

    Args:
        content_or_future: ChitchatResponder 返回的字符串或 coroutine (await 后得到字符串)
        request: OpenAI 兼容请求
        session_id: 会话 ID
        intent: 意图类型 ("short_query" 或 "off_topic")
    """
    completion_id = f"chatcmpl-{uuid.uuid4().hex[:24]}"
    created = int(time.time())

    span_name = f"agentinsight-researcher-{intent}"
    async with trace_agent(
        name=span_name,
        input={"session_id": session_id, "intent": intent},
        metadata={
            "session_id": session_id,
            "intent": intent,
            "mode": "chitchat_responder",
        },
        session_id=session_id,
    ):
        # await coroutine (ChitchatResponder.respond_xxx(stream=False) 返回 coroutine)
        content = (
            await content_or_future if not isinstance(content_or_future, str) else content_or_future
        )

    prompt_tokens = 0
    completion_tokens = len(content) // 4

    return ChatCompletionResponse(
        id=completion_id,
        created=created,
        model=request.model,
        choices=[
            {
                "index": 0,
                "message": {"role": "assistant", "content": content},
                "finish_reason": "stop",
            }
        ],
        usage={
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": prompt_tokens + completion_tokens,
        },
    )


# ========== 对话追问 (ChatAgentWithMemory) ==========


async def _stream_chat(
    initial_state: dict[str, Any],
    graph_config: dict[str, Any],
    request: ChatCompletionRequest,
    session_id: str,
    *,
    save_user_id: str | None = None,
    save_agent_id: str | None = None,
) -> Any:
    """流式 SSE 对话追问响应生成器.

    走 chat graph (单节点), 流式输出 AI 回答.
    用 trace_agent 包裹 graph.ainvoke 作为根 span.

    会话持久化: 流式结束后保存完整 assistant 响应到 chat_messages
    (save_user_id + save_agent_id 非空时触发).
    """
    completion_id = f"chatcmpl-{uuid.uuid4().hex[:24]}"
    created = int(time.time())

    def _sse_chunk(delta: dict[str, Any], finish_reason: str | None = None) -> str:
        """构造 SSE 数据帧."""
        chunk = {
            "id": completion_id,
            "object": "chat.completion.chunk",
            "created": created,
            "model": request.model,
            "choices": [
                {
                    "index": 0,
                    "delta": delta,
                    "finish_reason": finish_reason,
                }
            ],
        }
        return f"data: {orjson.dumps(chunk).decode('utf-8')}\n\n"

    # SSE 首块 (role)
    yield _sse_chunk({"role": "assistant"})

    # 收集流式内容 (用于会话持久化)
    collected_content: list[str] = []

    # 根 span 包裹对话追问
    async with trace_agent(
        name="agentinsight-researcher-chat",
        input={"query": initial_state["query"][:200], "session_id": session_id},
        metadata={
            "session_id": session_id,
            "intent": "chat",
            "user_id": initial_state.get("user_id"),
        },
        session_id=session_id,
        user_id=initial_state.get("user_id"),
    ):
        try:
            graph = await _get_chat_graph()

            async for event in graph.astream(
                initial_state, config=graph_config, stream_mode="updates"
            ):
                # event 格式: {node_name: delta_dict}
                for node_name, delta in event.items():
                    if not isinstance(delta, dict):
                        continue
                    if node_name == "chat":
                        # 提取 AIMessage 内容, 流式输出
                        msgs = delta.get("messages", [])
                        for msg in msgs:
                            if isinstance(msg, AIMessage) and msg.content:
                                msg_text = (
                                    msg.content
                                    if isinstance(msg.content, str)
                                    else str(msg.content)
                                )
                                collected_content.append(msg_text)
                                paragraphs = msg_text.split("\n\n")
                                for para in paragraphs:
                                    yield _sse_chunk({"content": para + "\n\n"})
                                    # 背压: 让出控制权, 允许消费者 (SSE writer) 刷出缓冲
                                    await asyncio.sleep(0.001)

        except Exception as e:
            logger.exception("对话追问流式执行失败")
            yield _sse_chunk({"content": f"\n\n**对话执行失败**: {str(e)[:200]}"})
            # 错误时 finish_reason="error"
            yield _sse_chunk({}, finish_reason="error")
            yield "data: [DONE]\n\n"
            return

    # 会话持久化: 保存 assistant 消息到 chat_messages
    # 必须在 yield [DONE] 之前, 否则 SSE 连接中断后 yield 之后的代码不执行, 导致历史消息丢失
    if save_user_id and save_agent_id and collected_content:
        full_response = "".join(collected_content)
        if full_response:
            await _persist_assistant_message(
                session_id,
                save_agent_id,
                save_user_id,
                full_response,
                metadata={"intent": "chat"},
            )

    # SSE 末块 (finish_reason)
    yield _sse_chunk({}, finish_reason="stop")
    yield "data: [DONE]\n\n"


async def _run_chat(
    initial_state: dict[str, Any],
    graph_config: dict[str, Any],
    request: ChatCompletionRequest,
    session_id: str,
) -> ChatCompletionResponse:
    """非流式对话追问执行.

    走 chat graph (单节点), 返回完整 AI 回答.
    trace_agent 根 span.
    """
    completion_id = f"chatcmpl-{uuid.uuid4().hex[:24]}"
    created = int(time.time())

    async with trace_agent(
        name="agentinsight-researcher-chat",
        input={"query": initial_state["query"][:200], "session_id": session_id},
        metadata={
            "session_id": session_id,
            "intent": "chat",
            "user_id": initial_state.get("user_id"),
        },
        session_id=session_id,
        user_id=initial_state.get("user_id"),
    ):
        settings = get_settings()
        try:
            graph = await _get_chat_graph()
            # graph.ainvoke 超时保护
            final_state = await asyncio.wait_for(
                graph.ainvoke(initial_state, config=graph_config),
                timeout=settings.graph_total_timeout,
            )

            # 从 messages 中提取最新 AIMessage 作为回答
            content = ""
            messages = final_state.get("messages", [])
            for msg in reversed(messages):
                if isinstance(msg, AIMessage) and msg.content:
                    content = msg.content if isinstance(msg.content, str) else str(msg.content)
                    break

            if not content:
                content = "(对话未生成响应)"

        except Exception as e:
            logger.exception("对话追问执行失败")
            content = f"对话执行失败: {str(e)[:500]}"

    prompt_tokens = len(initial_state["query"]) // 4
    completion_tokens = len(content) // 4

    return ChatCompletionResponse(
        id=completion_id,
        created=created,
        model=request.model,
        choices=[
            {
                "index": 0,
                "message": {"role": "assistant", "content": content},
                "finish_reason": "stop",
            }
        ],
        usage={
            "prompt_tokens": prompt_tokens,
            "completion_tokens": completion_tokens,
            "total_tokens": prompt_tokens + completion_tokens,
        },
    )


# ========== 文件上传 (用户需求 8) ==========

# 文件 magic number 签名表 (纯 Python, 无外部依赖, 避免新增 python-magic)
_FILE_MAGIC_SIGNATURES: dict[str, list[bytes]] = {
    "pdf": [b"%PDF"],
    # OOXML (docx/xlsx/pptx) 是 ZIP 格式: PK\x03\x04 (含文件) 或 PK\x05\x06 (空)
    "docx": [b"PK\x03\x04", b"PK\x05\x06"],
    "xlsx": [b"PK\x03\x04", b"PK\x05\x06"],
    "pptx": [b"PK\x03\x04", b"PK\x05\x06"],
    # 文本格式: UTF-8 可解码即通过 (无固定 magic)
    "md": [],
    "txt": [],
    "html": [],
    "csv": [],
}


def _validate_magic_number(data: bytes, ext: str) -> bool:
    """校验文件 magic number (防止恶意文件伪装扩展名).

    纯 Python 实现, 不依赖 python-magic / filetype 库 (避免不必要依赖).
    """
    signatures = _FILE_MAGIC_SIGNATURES.get(ext)
    if signatures is None:
        # 未注册扩展名 (已由扩展名白名单过滤), 跳过
        return True
    if not signatures:
        # 文本格式: 验证前 1KB 是否可 UTF-8 解码
        try:
            data[:1024].decode("utf-8")
            return True
        except UnicodeDecodeError:
            return False
    # 二进制格式: 校验 magic 前缀
    return any(data.startswith(sig) for sig in signatures)


@router.post("/files")
async def upload_file(
    file: UploadFile = File(...),  # noqa: B008 - FastAPI 标准模式
    authorization: str | None = Header(None),
) -> Any:
    """文件上传端点 (用户需求 8).

    上传文件作为研究数据源, 文件 ID 可在 /v1/chat/completions 的
    uploaded_files 字段引用.

    用户私有数据按 agent_id + user_id 隔离.
    安全约束 (大小/扩展名白名单/magic number).
    """
    settings = get_settings()
    user_id = get_request_user_id()
    agent_id = get_request_agent_id()

    # 校验扩展名 (白名单, 先于 I/O 校验)
    ext = Path(file.filename or "").suffix.lstrip(".").lower()
    if ext not in settings.allowed_extensions_list:
        raise HTTPException(
            status_code=415,
            detail=f"不支持的文件类型: .{ext}, 允许: {', '.join(settings.allowed_extensions_list)}",
        )

    # 优先用 file.size 早期校验大小 (UploadFile.size 由 Content-Length 注入)
    max_bytes = settings.max_upload_size_mb * 1024 * 1024
    if file.size is not None and file.size > max_bytes:
        raise HTTPException(
            status_code=413,
            detail=f"文件大小 {file.size / (1024 * 1024):.2f}MB 超过限制 {settings.max_upload_size_mb}MB",
        )

    # 分块流式读取 (避免全量读入内存), 同时校验大小 + magic number
    chunk_size = 1024 * 1024  # 1MB
    chunks: list[bytes] = []
    total_bytes = 0
    magic_checked = False
    while True:
        chunk = await file.read(chunk_size)
        if not chunk:
            break
        total_bytes += len(chunk)
        if total_bytes > max_bytes:
            raise HTTPException(
                status_code=413,
                detail=f"文件大小 {total_bytes / (1024 * 1024):.2f}MB 超过限制 {settings.max_upload_size_mb}MB",
            )
        # 首个 chunk 校验 magic number (防止恶意文件伪装扩展名)
        if not magic_checked:
            if not _validate_magic_number(chunk, ext):
                raise HTTPException(
                    status_code=415,
                    detail=f"文件内容与扩展名 .{ext} 不匹配 (magic number 校验失败)",
                )
            magic_checked = True
        chunks.append(chunk)

    contents = b"".join(chunks)
    size_mb = total_bytes / (1024 * 1024)

    # 生成文件 ID (agent_id:user_id:uuid 三级分键)
    file_id = f"{agent_id}:{user_id}:{uuid.uuid4().hex[:16]}"

    # 存储路径 (按 agent_id + user_id 隔离)
    upload_dir = Path(settings.upload_dir) / agent_id / user_id
    # 同步文件 I/O 经 asyncio.to_thread 包裹, 避免阻塞事件循环
    await asyncio.to_thread(upload_dir.mkdir, parents=True, exist_ok=True)
    save_path = upload_dir / f"{file_id.split(':')[-2]}_{file_id.split(':')[-1]}.{ext}"

    # 写入文件 (asyncio.to_thread 包裹同步 write_bytes)
    await asyncio.to_thread(save_path.write_bytes, contents)

    logger.info(
        "文件上传成功: file_id=%s, filename=%s, size=%.2fMB, user=%s",
        file_id,
        file.filename,
        size_mb,
        user_id,
    )

    return JSONResponse(
        status_code=201,
        content={
            "file_id": file_id,
            "filename": file.filename,
            "size_bytes": total_bytes,
            "size_mb": round(size_mb, 4),
            "extension": ext,
            "uploaded_at": int(time.time()),
        },
    )


async def _load_uploaded_files_context(
    file_ids: list[str], user_id: str, agent_id: str
) -> list[str]:
    """加载已上传文件内容作为研究上下文.

    按 agent_id + user_id 隔离, 禁止跨用户访问.
    所有同步文件 I/O (exists/glob/read_text/第三方库 open) 经 asyncio.to_thread
    包裹, 避免阻塞事件循环. _extract_file_content 内部含 fitz/Document/openpyxl/pptx
    等同步库调用, 整体托管到线程执行.
    """
    settings = get_settings()
    contexts: list[str] = []

    for file_id in file_ids:
        try:
            # 校验 file_id 前缀归属当前 agent+user (安全: 禁止跨用户)
            parts = file_id.split(":")
            if len(parts) != 3:
                continue
            fid_agent, fid_user, fid_uuid = parts
            if fid_agent != agent_id or fid_user != user_id:
                logger.warning(
                    "拒绝跨用户文件访问: file_id=%s, agent=%s, user=%s",
                    file_id,
                    agent_id,
                    user_id,
                )
                continue

            # 查找文件 (按 uuid 前缀匹配)
            upload_dir = Path(settings.upload_dir) / agent_id / user_id
            # exists/glob 同步 I/O 经 asyncio.to_thread 包裹
            if not await asyncio.to_thread(upload_dir.exists):
                continue
            # B023: 显式绑定循环变量到 lambda 默认参数, 避免闭包延迟绑定
            _upload_dir = upload_dir
            _fid_uuid = fid_uuid

            def _glob_files(d: Path = _upload_dir, u: str = _fid_uuid) -> list[Path]:
                return list(d.glob(f"{u}_*"))

            matches = await asyncio.to_thread(_glob_files)
            if not matches:
                continue

            file_path = matches[0]
            # _extract_file_content 含 read_text/fitz.open/Document 等同步 I/O,
            # 整体经 asyncio.to_thread 包裹, 避免阻塞事件循环
            content = await asyncio.to_thread(
                _extract_file_content,
                file_path,
                file_path.suffix.lstrip(".").lower(),
            )
            if content:
                contexts.append(f"=== 用户上传文件: {file_path.name} ===\n{content[:8000]}")
        except Exception as e:  # noqa: BLE001
            logger.warning("加载文件 %s 失败: %s", file_id, e)

    return contexts


def _extract_file_content(file_path: Path, ext: str) -> str:
    """提取文件文本内容 (按扩展名路由).

    支持: pdf, docx, md, txt, html, csv (用户需求 8).
    """
    try:
        if ext in ("txt", "md", "csv"):
            return str(file_path.read_text(encoding="utf-8", errors="ignore"))

        if ext == "pdf":
            try:
                from pypdf import PdfReader

                reader = PdfReader(str(file_path))
                text = "\n".join(page.extract_text() or "" for page in reader.pages)
                return str(text)
            except ImportError:
                logger.warning("pypdf 未安装, 跳过 PDF 文本提取")
                return ""

        if ext == "docx":
            try:
                from docx import Document

                doc = Document(str(file_path))
                return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            except ImportError:
                logger.warning("python-docx 未安装, 跳过 DOCX 文本提取")
                return ""

        if ext == "html":
            try:
                from bs4 import BeautifulSoup

                html = file_path.read_text(encoding="utf-8", errors="ignore")
                soup = BeautifulSoup(html, "html.parser")
                return str(soup.get_text(separator="\n", strip=True))
            except ImportError:
                return str(file_path.read_text(encoding="utf-8", errors="ignore"))

        if ext == "xlsx":
            try:
                from openpyxl import load_workbook

                wb = load_workbook(str(file_path), read_only=True)
                text_parts: list[str] = []
                for ws in wb.worksheets:
                    for row in ws.iter_rows(values_only=True):
                        text_parts.append(",".join(str(c) if c is not None else "" for c in row))
                wb.close()
                return "\n".join(text_parts)
            except ImportError:
                logger.warning("openpyxl 未安装, 跳过 XLSX 文本提取")
                return ""

        if ext == "pptx":
            try:
                from pptx import Presentation

                prs = Presentation(str(file_path))
                pptx_parts: list[str] = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            pptx_parts.append(shape.text)
                return "\n".join(pptx_parts)
            except ImportError:
                logger.warning("python-pptx 未安装, 跳过 PPTX 文本提取")
                return ""

    except Exception as e:  # noqa: BLE001
        logger.warning("提取文件 %s 内容失败: %s", file_path.name, e)

    return ""


@router.get("/models")
async def list_models() -> Any:
    """OpenAI 兼容 /v1/models 端点."""
    return {
        "object": "list",
        "data": [
            {
                "id": "agentinsight-researcher",
                "object": "model",
                "created": int(time.time()),
                "owned_by": "agentinsight",
            }
        ],
    }


# ========== 人在回路反馈 (Human-in-the-loop) ==========


class FeedbackRequest(BaseModel):
    """人在回路反馈请求.

    /v1/feedback 为允许调用的端点 (人在回路反馈通道).
    用户通过此端点提交对研究计划/大纲的审核反馈, 解决 HumanAgent 等待的 Future.
    """

    session_id: str = Field(..., description="会话 ID (thread_id), 与研究请求一致")
    feedback: str = Field(
        ...,
        description=(
            "审核反馈内容; 空字符串或 approve/accept/通过 等关键词表示接受, "
            "其他内容视为修订意见 (回 agent_creator 重新生成角色)."
        ),
    )


@router.post("/feedback")
async def submit_feedback(request: FeedbackRequest) -> Any:
    """提交人在回路审核反馈.

    /v1/feedback 为允许调用的端点.
    HumanAgent 在 human 节点通过 FeedbackQueue.wait_feedback() 阻塞等待,
    此端点调用 FeedbackQueue.put_feedback() 提交反馈, 解决等待.

    Returns:
        200: 反馈已提交
        404: 无待处理的反馈请求 (HumanAgent 未在等待)
    """
    from src.api.feedback_queue import get_feedback_queue

    feedback_queue = get_feedback_queue()
    ok = feedback_queue.put_feedback(request.session_id, request.feedback)
    if not ok:
        raise HTTPException(
            status_code=404,
            detail="无待处理的反馈请求 (session_id 可能无效或反馈已提交)",
        )
    return JSONResponse(
        status_code=200,
        content={
            "session_id": request.session_id,
            "submitted": True,
            "submitted_at": int(time.time()),
        },
    )


# ========== 报告下载端点 (多格式输出支持) ==========


@router.get("/reports/session/{session_id}")
async def list_session_reports(
    session_id: str,
    authorization: str | None = Header(None),
) -> list[dict[str, Any]]:
    """列出指定 session 的所有报告.

    一个 session 可以生成多个报告, 返回按 created_at DESC 排序的列表.
    每项含: report_id, session_id, query, report_format, created_at, updated_at.

    数据按 agent_id + user_id 隔离, 此处复用请求上下文 user_id.
    """
    from src.memory.report_store import get_report_store

    user_id = get_request_user_id() or "anonymous"
    store = get_report_store()
    reports = await store.list_reports(
        session_id=session_id,
        user_id=user_id,
        limit=50,
    )
    # 仅返回列表展示需要的字段 (不含 report_md 全文, 减少传输)
    return [
        {
            "report_id": r.get("report_id"),
            "session_id": r.get("session_id"),
            "query": (r.get("query", "") or "")[:200],
            "report_format": r.get("report_format", "markdown"),
            "agent_role": r.get("agent_role"),
            "created_at": r.get("created_at"),
            "updated_at": r.get("updated_at"),
        }
        for r in reports
    ]


@router.get("/reports/{report_id}/download")
async def download_report(
    report_id: str,
    format: str = "markdown",
    authorization: str | None = Header(None),
) -> Response:
    """下载研究报告文件 (按 report_id, 支持多格式实时转换).

    支持 format: markdown / html / pdf / docx / json
    数据按 agent_id + user_id 隔离.

    向后兼容: 若 report_id 未匹配到记录, 尝试将其作为 session_id 查询最新报告
    (兼容字段, 响应头 X-Deprecated 提示调用方使用 report_id).
    """
    import aiofiles
    import aiofiles.os

    from src.memory.report_store import get_report_store
    from src.skills.researcher.publisher import Publisher

    user_id = get_request_user_id() or "anonymous"
    store = get_report_store()

    # report_id 应为 UUID, 不符合格式时跳过 get_report (避免 asyncpg $1::uuid 抛 DataError → 500)
    # 直接走 session_id 兼容分支 (session_id 为 VARCHAR(64), 接受任意字符串)
    report: dict[str, Any] | None = None
    try:
        uuid.UUID(report_id)
    except ValueError:
        report = None
    else:
        report = await store.get_report(report_id)

    deprecated_fallback = False
    if not report:
        # 向后兼容: 调用方仍传 session_id (旧逻辑), 取该 session 最新报告下载
        reports = await store.list_reports(
            session_id=report_id,
            user_id=user_id,
            limit=1,
        )
        if reports:
            report = reports[0]
            deprecated_fallback = True
            logger.warning(
                "下载端点收到 session_id=%s (兼容字段), 建议使用 report_id=%s",
                report_id,
                report.get("report_id"),
            )

    if not report:
        raise HTTPException(status_code=404, detail="报告不存在")

    # 数据隔离校验: report.user_id 必须匹配当前 user_id (匿名用户跳过)
    report_user_id = report.get("user_id")
    if report_user_id and report_user_id != user_id:
        raise HTTPException(status_code=403, detail="无权访问该报告")

    actual_report_id = report.get("report_id") or report_id
    content = report.get("report_md", "")
    title = (report.get("query", "") or "研究报告")[:100]
    sources = report.get("sources", []) if isinstance(report.get("sources"), list) else []
    agent_role = report.get("agent_role", "") or ""

    publisher = Publisher()
    # 公共响应头: 含 deprecation 标记 (兼容旧 session_id 调用时)
    extra_headers: dict[str, str] = {}
    if deprecated_fallback:
        extra_headers["X-Deprecated"] = (
            "true - 使用 session_id 调用已弃用, 请改用 report_id (参考 /v1/reports/session/{session_id})"
        )
        extra_headers["Link"] = (
            f'</v1/reports/{actual_report_id}/download>; rel="successor-version"'
        )

    if format == "markdown":
        return Response(
            content=content.encode("utf-8"),
            media_type="text/markdown",
            headers={
                "Content-Disposition": f"attachment; filename=report_{actual_report_id}.md",
                **extra_headers,
            },
        )
    elif format == "html":
        html = publisher._md_to_html(content)
        return Response(
            content=html.encode("utf-8"),
            media_type="text/html",
            headers={
                "Content-Disposition": f"attachment; filename=report_{actual_report_id}.html",
                **extra_headers,
            },
        )
    elif format == "pdf":
        pdf_path = await publisher._md_to_pdf(content, actual_report_id)
        if await aiofiles.os.path.exists(pdf_path):
            async with aiofiles.open(pdf_path, "rb") as f:
                pdf_content = await f.read()
            return Response(
                content=pdf_content,
                media_type="application/pdf",
                headers={
                    "Content-Disposition": f"attachment; filename=report_{actual_report_id}.pdf",
                    **extra_headers,
                },
            )
        raise HTTPException(status_code=404, detail="PDF 生成失败")
    elif format == "docx":
        docx_bytes = publisher._to_docx(content, title=title)
        return Response(
            content=docx_bytes,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={
                "Content-Disposition": f"attachment; filename=report_{actual_report_id}.docx",
                **extra_headers,
            },
        )
    elif format == "json":
        json_str = publisher._to_json(
            content,
            title=title,
            sources=sources,
            agent_role_server=agent_role,
        )
        return Response(
            content=json_str.encode("utf-8"),
            media_type="application/json",
            headers={
                "Content-Disposition": f"attachment; filename=report_{actual_report_id}.json",
                **extra_headers,
            },
        )
    else:
        raise HTTPException(status_code=400, detail=f"不支持的格式: {format}")
