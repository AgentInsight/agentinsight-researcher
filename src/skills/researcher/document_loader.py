"""文档加载模块.

统一封装多种数据源 (URL/本地文件/Azure Blob) 的文档加载, 返回 list[Document].

设计要点:
- DocumentLoader 基类: async load(source) -> list[Document], 单一职责.
- OnlineDocumentLoader: 复用现有 scrapers (scrapers/__init__.py scrape_with_fallback).
- LocalDocumentLoader: 按扩展名路由 (PDF/Word/Excel/Markdown/TXT/HTML/CSV/PPTX).
- AzureBlobLoader: 可选, 需 azure-storage-blob (用 try/except import).
- 工厂函数 get_document_loader(source): 按 source 形态自动路由.

接入主流程 (用户文件上传 RAG 链路):
    routes.py upload_file 端点应调用本模块提取文档内容, 再调用
    EmbeddingsClient.embed_and_index 索引到 Qdrant 用户私有 namespace:
        loader = get_document_loader(file_path, settings)
        docs = await loader.load(file_path)
        texts = [d.page_content for d in docs]
        await embeddings.embed_and_index(
            texts,
            namespace=f"{agent_id}:{user_id}",
            user_id=user_id,
            session_id=session_id,
        )
    注: routes.py 的 upload_file 端点接入由调用方负责, 本模块仅保证工厂可用.

依赖说明 (不在 requirements.txt 中的可选依赖):
- PyMuPDF (fitz): 已在 requirements.txt (PDF 解析).
- python-docx / openpyxl / python-pptx: 已在 requirements.txt.
- azure-storage-blob: 不在 requirements.txt, 需手动 pip install azure-storage-blob.
- markitdown: 已在 requirements.txt (Office 文档兜底).

AGENTS.md 第 9 章: 工具调用必须经 trace_xxx span 包裹; 本模块属文档加载,
非 MCP 工具, 不强制 trace_tool; 调用方如需可在 chain span 内包裹.
"""

from __future__ import annotations

import asyncio
import logging
import os
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from typing import Any

from src.config.settings import Settings, get_settings

logger = logging.getLogger(__name__)


# ========== Document 数据类 (薄封装, 避免引入额外耦合) ==========


@dataclass
class Document:
    """加载后的文档对象.

    自研薄封装, 避免引入额外耦合.
    """

    page_content: str
    metadata: dict[str, Any] = field(default_factory=dict)


# ========== DocumentLoader 基类 ==========


class DocumentLoader(ABC):
    """文档加载器基类.

    所有 loader 共享 load(source) -> list[Document] 规约.
    """

    def __init__(self, settings: Settings | None = None) -> None:
        self.settings = settings or get_settings()

    @abstractmethod
    async def load(self, source: str) -> list[Document]:
        """加载文档, 返回 list[Document].

        Args:
            source: 数据源标识 (URL / 本地路径 / Blob URI).
        """
        raise NotImplementedError


# ========== OnlineDocumentLoader (URL → scrapers) ==========


class OnlineDocumentLoader(DocumentLoader):
    """从 URL 加载文档 (复用现有 scrapers).

    走 src.skills.researcher.scrapers.scrape_with_fallback 降级链.
    """

    async def load(self, source: str) -> list[Document]:
        """从 URL 抓取并转为 Document."""
        if not source.startswith(("http://", "https://")):
            return []

        from src.skills.researcher.scrapers import scrape_with_fallback

        try:
            result = await scrape_with_fallback(
                source,
                enable_fallback=True,
                min_content_length=100,
            )
        except Exception as e:  # noqa: BLE001
            logger.warning("OnlineDocumentLoader 抓取失败 %s: %s", source, e)
            return []

        content = result.get("content") or ""
        if not content:
            return []

        return [
            Document(
                page_content=content,
                metadata={
                    "source": source,
                    "title": result.get("title", ""),
                    "content_type": result.get("content_type", "html"),
                    "image_urls": result.get("image_urls", []),
                },
            )
        ]


# ========== LocalDocumentLoader (本地文件) ==========


class LocalDocumentLoader(DocumentLoader):
    """从本地文件路径加载文档.

    按扩展名路由:
    - .pdf       → PyMuPDF (fitz)
    - .docx      → python-docx
    - .xlsx/.xls → openpyxl
    - .pptx      → python-pptx
    - .md/.txt/.csv/.html → 直接读取 (HTML 用 BeautifulSoup)
    """

    async def load(self, source: str) -> list[Document]:
        """加载本地文件, 返回 list[Document]."""
        if not source or source.startswith(("http://", "https://")):
            return []

        # 检查文件存在 (同步 IO 用 to_thread 包装避免 ASYNC230)
        def _exists(path: str) -> bool:
            return os.path.exists(path) and os.path.isfile(path)

        if not await asyncio.to_thread(_exists, source):
            logger.warning("LocalDocumentLoader 文件不存在: %s", source)
            return []

        ext = self._get_extension(source)
        try:
            content = await asyncio.to_thread(self._extract, source, ext)
        except Exception as e:  # noqa: BLE001
            logger.warning("LocalDocumentLoader 提取失败 %s: %s", source, e)
            return []

        if not content:
            return []

        return [
            Document(
                page_content=content,
                metadata={
                    "source": source,
                    "file_type": ext,
                    "file_name": os.path.basename(source),
                },
            )
        ]

    @staticmethod
    def _get_extension(path: str) -> str:
        """提取文件扩展名 (小写, 不含点)."""
        _, ext = os.path.splitext(path)
        return ext.lstrip(".").lower()

    def _extract(self, path: str, ext: str) -> str:
        """按扩展名路由提取文本 (同步, 由 load() 用 to_thread 包装).

        可选依赖未安装时记录告警并返回空串, 不抛 ImportError.
        """
        if ext in ("txt", "md", "csv"):
            return self._read_text(path)

        if ext == "pdf":
            return self._extract_pdf(path)

        if ext == "docx":
            return self._extract_docx(path)

        if ext in ("xlsx", "xls"):
            return self._extract_xlsx(path)

        if ext == "pptx":
            return self._extract_pptx(path)

        if ext == "html":
            return self._extract_html(path)

        # 未知扩展名: 尝试按文本读取
        logger.info("LocalDocumentLoader 未知扩展名 %s, 尝试按文本读取: %s", ext, path)
        return self._read_text(path)

    @staticmethod
    def _read_text(path: str) -> str:
        """直接读取文本文件."""
        with open(path, encoding="utf-8", errors="ignore") as f:
            return f.read()

    @staticmethod
    def _extract_pdf(path: str) -> str:
        """用 PyMuPDF (fitz) 提取 PDF 文本."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.warning("PyMuPDF (fitz) 未安装, 无法解析 PDF: %s", path)
            return ""
        doc = fitz.open(path)
        try:
            return "\n\n".join(page.get_text() for page in doc)
        finally:
            doc.close()

    @staticmethod
    def _extract_docx(path: str) -> str:
        """用 python-docx 提取 DOCX 文本."""
        try:
            from docx import Document as DocxDocument
        except ImportError:
            logger.warning("python-docx 未安装, 无法解析 DOCX: %s", path)
            return ""
        doc = DocxDocument(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    @staticmethod
    def _extract_xlsx(path: str) -> str:
        """用 openpyxl 提取 XLSX 文本 (按行 CSV 化)."""
        try:
            from openpyxl import load_workbook
        except ImportError:
            logger.warning("openpyxl 未安装, 无法解析 XLSX: %s", path)
            return ""
        wb = load_workbook(path, read_only=True)
        try:
            parts: list[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    parts.append(",".join(str(c) if c is not None else "" for c in row))
            return "\n".join(parts)
        finally:
            wb.close()

    @staticmethod
    def _extract_pptx(path: str) -> str:
        """用 python-pptx 提取 PPTX 文本."""
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("python-pptx 未安装, 无法解析 PPTX: %s", path)
            return ""
        prs = Presentation(path)
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts)

    @staticmethod
    def _extract_html(path: str) -> str:
        """提取 HTML 文本 (优先 BeautifulSoup, 降级纯文本)."""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            logger.warning("beautifulsoup4 未安装, HTML 按纯文本读取: %s", path)
            with open(path, encoding="utf-8", errors="ignore") as f:
                return f.read()
        with open(path, encoding="utf-8", errors="ignore") as f:
            html = f.read()
        soup = BeautifulSoup(html, "html.parser")
        return str(soup.get_text(separator="\n", strip=True))


# ========== AzureBlobLoader (Azure Blob Storage, 可选依赖) ==========


class AzureBlobLoader(DocumentLoader):
    """从 Azure Blob Storage 加载文档 (可选, 需 azure-storage-blob).

    支持 source 形态:
    - azure://container/blob
    - https://<account>.blob.core.windows.net/<container>/<blob>

    azure-storage-blob 未安装时记录告警并返回空列表.
    """

    async def load(self, source: str) -> list[Document]:
        """从 Azure Blob 加载文档, 返回 list[Document]."""
        try:
            from azure.storage.blob import BlobServiceClient
        except ImportError:
            logger.warning("azure-storage-blob 未安装, AzureBlobLoader 无法加载: %s", source)
            return []

        try:
            container_name, blob_name = self._parse_source(source)
        except ValueError as e:
            logger.warning("AzureBlobLoader source 解析失败 %s: %s", source, e)
            return []

        try:
            conn_str = self.settings.azure_storage_connection_string
            if not conn_str:
                logger.warning(
                    "azure_storage_connection_string 未配置, AzureBlobLoader 无法加载: %s",
                    source,
                )
                return []
            client = BlobServiceClient.from_connection_string(conn_str)
            blob_client = client.get_blob_client(container=container_name, blob=blob_name)
            downloader = await asyncio.to_thread(blob_client.download_blob)
            content = await asyncio.to_thread(downloader.readall)
            text = content.decode("utf-8", errors="ignore")
            return [
                Document(
                    page_content=text,
                    metadata={
                        "source": source,
                        "container": container_name,
                        "blob": blob_name,
                        "content_type": "azure_blob",
                    },
                )
            ]
        except Exception as e:  # noqa: BLE001
            logger.warning("AzureBlobLoader 加载失败 %s: %s", source, e)
            return []

    @staticmethod
    def _parse_source(source: str) -> tuple[str, str]:
        """解析 source 为 (container, blob).

        支持:
        - azure://container/blob
        - https://<account>.blob.core.windows.net/<container>/<blob>
        """
        if source.startswith("azure://"):
            path = source[len("azure://") :]
        elif "blob.core.windows.net" in source:
            # https://<account>.blob.core.windows.net/<container>/<blob>
            after_netloc = source.split("blob.core.windows.net/", 1)[-1]
            path = after_netloc
        else:
            raise ValueError(f"不支持的 Azure Blob source 形态: {source}")

        parts = path.split("/", 1)
        if len(parts) != 2 or not all(parts):
            raise ValueError(f"Azure Blob source 缺少 container/blob: {source}")
        return parts[0], parts[1]


# ========== 工厂函数 ==========


def get_document_loader(
    source: str,
    settings: Settings | None = None,
) -> DocumentLoader:
    """根据 source 形态自动选择 DocumentLoader.

    工厂路由逻辑:
    - source 以 http://|https:// 开头 → OnlineDocumentLoader
    - source 以 azure:// 或 Azure Blob URL 形式 → AzureBlobLoader
    - source 以本地路径存在 (或非 URL) → LocalDocumentLoader

    Args:
        source: 数据源 (URL / 本地路径 / Blob URI).
        settings: 可选 Settings 注入 (默认 get_settings()).
    """
    settings = settings or get_settings()

    if source.startswith(("http://", "https://")):
        # Azure Blob URL 简单识别 (blob.core.windows.net)
        if "blob.core.windows.net" in source:
            return AzureBlobLoader(settings)
        return OnlineDocumentLoader(settings)

    # azure://container/blob 协议
    if source.startswith("azure://"):
        return AzureBlobLoader(settings)

    # 默认按本地路径处理
    return LocalDocumentLoader(settings)
